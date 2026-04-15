"""
Netgear NRH Report Generator — Unified CLI tool.

Merges the functionality of netgear_report.py, netgear_report_config.py,
and netgear_report_interactive.py into a single configurable script.

Usage:
    # With config file
    python -m cli.report_cli --config report_config.json

    # Interactive mode
    python -m cli.report_cli --interactive --excel data.xlsm --template report.pptx

    # Direct mode with inline mappings
    python -m cli.report_cli --excel data.xlsm --template report.pptx \\
        --output result.pptx --map "Metric DUT vs REF#1:8:worksheet" --map "BI:9:chartsheet"
"""
import os
import sys
import json
import argparse

# Add project root to path for imports
sys.path.insert(0, os.path.join(os.path.dirname(__file__), ".."))

from app.config import logger, DEFAULT_IMAGE_LAYOUT


def parse_args():
    p = argparse.ArgumentParser(description="Netgear NRH Report Generator CLI")
    p.add_argument("--excel", help="Excel file path")
    p.add_argument("--template", help="PPT template path")
    p.add_argument("--output", help="Output PPT path")
    p.add_argument("--config", help="JSON config file with mappings")
    p.add_argument("--interactive", action="store_true", help="Interactive mode")
    p.add_argument(
        "--map",
        action="append",
        default=[],
        help="Mapping in format 'name:page:type' (can be repeated)",
    )
    p.add_argument("--img-left", type=float, default=DEFAULT_IMAGE_LAYOUT["left"])
    p.add_argument("--img-top", type=float, default=DEFAULT_IMAGE_LAYOUT["top"])
    p.add_argument("--img-width", type=float, default=DEFAULT_IMAGE_LAYOUT["width"])
    p.add_argument("--img-height", type=float, default=DEFAULT_IMAGE_LAYOUT["height"])
    return p.parse_args()


def load_config(config_path: str) -> dict:
    """Load configuration from a JSON file."""
    with open(config_path, "r", encoding="utf-8") as f:
        return json.load(f)


def list_available_items(excel_path: str):
    """List all available worksheets and chart sheets."""
    from app.services.excel_service import ExcelCOM

    with ExcelCOM() as (excel_app, _):
        workbook = excel_app.Workbooks.Open(os.path.abspath(excel_path))
        worksheets = [sheet.Name for sheet in workbook.Worksheets]
        chartsheets = []
        try:
            for chart in workbook.Charts:
                chartsheets.append(chart.Name)
        except Exception:
            pass
        workbook.Close(SaveChanges=False)

    return worksheets, chartsheets


def interactive_mode(excel_path: str, template_path: str, output_path: str, args):
    """Run in interactive mode, prompting user for mappings."""
    from pptx import Presentation

    print("\n" + "=" * 60)
    print("  Netgear NRH Report Generator — Interactive Mode")
    print("=" * 60)

    print(f"\n  Excel:    {os.path.basename(excel_path)}")
    print(f"  Template: {os.path.basename(template_path)}")
    print(f"  Output:   {os.path.basename(output_path)}")

    print("\n  Scanning Excel file...")
    worksheets, chartsheets = list_available_items(excel_path)

    print("\n" + "-" * 60)
    print("  Available items:")
    print("-" * 60)

    all_items = []
    idx = 1
    print("\n  [Worksheets]")
    for ws in worksheets:
        print(f"    {idx}. {ws}")
        all_items.append((ws, "worksheet"))
        idx += 1

    print("\n  [Chart Sheets]")
    for cs in chartsheets:
        print(f"    {idx}. {cs}")
        all_items.append((cs, "chartsheet"))
        idx += 1

    prs = Presentation(template_path)
    total_slides = len(prs.slides)
    print(f"\n  PPT has {total_slides} slides")

    print("\n" + "=" * 60)
    print("  Enter selections:  <item#>,<page#>")
    print("  Type 'done' when finished, 'quit' to cancel")
    print("-" * 60)

    selections = []
    while True:
        user_input = input("\n  > ").strip().lower()
        if user_input == "done":
            break
        elif user_input == "quit":
            print("  Cancelled.")
            return

        try:
            parts = user_input.split(",")
            item_num = int(parts[0].strip())
            page_num = int(parts[1].strip())

            if item_num < 1 or item_num > len(all_items):
                print(f"    Invalid item. Must be 1-{len(all_items)}")
                continue
            if page_num < 1 or page_num > total_slides:
                print(f"    Invalid page. Must be 1-{total_slides}")
                continue

            name, item_type = all_items[item_num - 1]
            selections.append({"name": name, "type": item_type, "page": page_num})
            print(f"    Added: {name} -> Page {page_num}")

        except (ValueError, IndexError):
            print("    Invalid format. Use: <item>,<page>  e.g. 4,8")

    if not selections:
        print("\n  No selections made. Exiting.")
        return

    print("\n  Your selections:")
    for sel in selections:
        print(f"    {sel['name']} -> Page {sel['page']}")

    confirm = input("\n  Proceed? (y/n): ").strip().lower()
    if confirm != "y":
        print("  Cancelled.")
        return

    run_generation(excel_path, template_path, output_path, selections, args)


def run_generation(excel_path, template_path, output_path, mappings, args):
    """Execute the actual extraction and insertion."""
    from app.services.excel_service import ExcelCOM, capture_item
    from pptx import Presentation
    from pptx.util import Inches

    temp_dir = os.path.join(os.path.dirname(output_path), "_temp_charts")
    os.makedirs(temp_dir, exist_ok=True)

    # Step 1: Extract
    print("\n" + "=" * 60)
    print("  Step 1: Extracting from Excel")
    print("=" * 60)

    extracted = {}
    with ExcelCOM() as (excel_app, _):
        workbook = excel_app.Workbooks.Open(os.path.abspath(excel_path))

        for sel in mappings:
            name = sel["name"]
            safe_name = name.replace(" ", "_").replace("#", "_").replace("/", "_")
            img_path = os.path.join(temp_dir, f"{safe_name}.png")

            print(f"\n  Extracting: {name}")
            if capture_item(excel_app, workbook, name, sel["type"], img_path):
                extracted[name] = img_path
                print(f"    OK ({os.path.getsize(img_path)} bytes)")
            else:
                print(f"    FAILED")

        workbook.Close(SaveChanges=False)

    # Step 2: Insert into PPT
    print("\n" + "=" * 60)
    print("  Step 2: Inserting into PowerPoint")
    print("=" * 60)

    prs = Presentation(template_path)

    for sel in mappings:
        name = sel["name"]
        page = sel["page"]
        slide_idx = page - 1

        if name not in extracted:
            print(f"\n  SKIP: {name} — no image")
            continue

        if slide_idx >= len(prs.slides):
            print(f"\n  SKIP: Page {page} doesn't exist")
            continue

        print(f"\n  {name} -> Page {page}")
        slide = prs.slides[slide_idx]

        slide.shapes.add_picture(
            extracted[name],
            Inches(args.img_left),
            Inches(args.img_top),
            width=Inches(args.img_width),
            height=Inches(args.img_height),
        )
        print(f"    OK")

    prs.save(output_path)

    # Cleanup temp
    import shutil
    shutil.rmtree(temp_dir, ignore_errors=True)

    print("\n" + "=" * 60)
    print(f"  Done! Output: {output_path}")
    print("=" * 60)


def main():
    args = parse_args()

    # Load config file if provided
    if args.config:
        cfg = load_config(args.config)
        excel_path = cfg.get("excel_file", args.excel)
        template_path = cfg.get("ppt_template", args.template)
        output_path = cfg.get("output_ppt", args.output)
        mappings = [
            {"name": m[0], "page": m[1], "type": m[2]}
            for m in cfg.get("mappings", [])
        ]

        if not excel_path or not template_path or not output_path:
            print("Error: config must include excel_file, ppt_template, output_ppt")
            return

        run_generation(excel_path, template_path, output_path, mappings, args)
        return

    # Validate required paths
    if not args.excel or not args.template:
        print("Error: --excel and --template are required (or use --config)")
        return

    output = args.output or os.path.splitext(args.template)[0] + "_Generated.pptx"

    if args.interactive:
        interactive_mode(args.excel, args.template, output, args)
        return

    # Direct mode with --map flags
    if not args.map:
        print("Error: provide --map flags or use --interactive / --config")
        return

    mappings = []
    for m in args.map:
        parts = m.split(":")
        if len(parts) != 3:
            print(f"Error: invalid mapping format '{m}'. Use 'name:page:type'")
            return
        mappings.append({"name": parts[0], "page": int(parts[1]), "type": parts[2]})

    run_generation(args.excel, args.template, output, mappings, args)


if __name__ == "__main__":
    main()
