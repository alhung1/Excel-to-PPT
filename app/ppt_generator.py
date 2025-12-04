"""
PowerPoint Generator
Creates presentations from template and fills with chart images
"""
import os
from typing import List, Dict, Optional, Tuple
from pptx import Presentation
from pptx.util import Inches, Pt


class PPTGenerator:
    """Generate PowerPoint presentations with chart images"""
    
    def __init__(self, template_path: Optional[str] = None):
        """
        Initialize generator with optional template
        
        Args:
            template_path: Path to PowerPoint template (.pptx)
        """
        if template_path and os.path.exists(template_path):
            self.prs = Presentation(template_path)
            self.template_path = template_path
        else:
            self.prs = Presentation()
            self.template_path = None
            # Set default slide size to widescreen 16:9
            self.prs.slide_width = Inches(13.333)
            self.prs.slide_height = Inches(7.5)
    
    def add_title_slide(self, title: str, subtitle: str = "") -> None:
        """Add a title slide"""
        slide_layout = self.prs.slide_layouts[0]  # Title Slide layout
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Set title
        if slide.shapes.title:
            slide.shapes.title.text = title
        
        # Set subtitle if placeholder exists
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 1:  # Subtitle placeholder
                shape.text = subtitle
                break
    
    def add_chart_slide(
        self, 
        chart_image_path: str,
        title: str = "",
        position: Optional[Tuple[float, float, float, float]] = None
    ) -> None:
        """
        Add a slide with a chart image
        
        Args:
            chart_image_path: Path to the chart image
            title: Optional slide title
            position: Optional (left, top, width, height) in inches
        """
        # Use blank or content layout
        layout_idx = 6 if len(self.prs.slide_layouts) > 6 else 5
        try:
            slide_layout = self.prs.slide_layouts[layout_idx]
        except:
            slide_layout = self.prs.slide_layouts[-1]
            
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Add title if provided
        if title:
            if slide.shapes.title:
                slide.shapes.title.text = title
            else:
                # Add text box for title
                left = Inches(0.5)
                top = Inches(0.3)
                width = Inches(12)
                height = Inches(0.6)
                txBox = slide.shapes.add_textbox(left, top, width, height)
                tf = txBox.text_frame
                tf.text = title
                tf.paragraphs[0].font.size = Pt(28)
                tf.paragraphs[0].font.bold = True
        
        # Calculate image position
        if position:
            left, top, width, height = [Inches(x) for x in position]
        else:
            # Center the image with some padding
            left = Inches(0.75)
            top = Inches(1.5) if title else Inches(0.75)
            width = Inches(11.8)
            height = Inches(5.5) if title else Inches(6.5)
        
        # Add chart image
        if os.path.exists(chart_image_path):
            slide.shapes.add_picture(chart_image_path, left, top, width=width)
    
    def add_multiple_charts_slide(
        self,
        chart_images: List[str],
        title: str = "",
        layout: str = "grid"  # "grid", "horizontal", "vertical"
    ) -> None:
        """
        Add a slide with multiple chart images
        
        Args:
            chart_images: List of paths to chart images
            title: Optional slide title
            layout: How to arrange charts
        """
        layout_idx = 6 if len(self.prs.slide_layouts) > 6 else 5
        try:
            slide_layout = self.prs.slide_layouts[layout_idx]
        except:
            slide_layout = self.prs.slide_layouts[-1]
            
        slide = self.prs.slides.add_slide(slide_layout)
        
        # Add title
        if title:
            left = Inches(0.5)
            top = Inches(0.3)
            width = Inches(12)
            height = Inches(0.6)
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.text = title
            tf.paragraphs[0].font.size = Pt(24)
            tf.paragraphs[0].font.bold = True
        
        # Calculate positions based on layout
        content_top = 1.3 if title else 0.5
        content_height = 6.0 if title else 6.8
        num_charts = len(chart_images)
        
        positions = self._calculate_positions(num_charts, layout, content_top, content_height)
        
        # Add chart images
        for i, (chart_path, pos) in enumerate(zip(chart_images, positions)):
            if os.path.exists(chart_path):
                left, top, width, height = [Inches(x) for x in pos]
                slide.shapes.add_picture(chart_path, left, top, width=width)
    
    def _calculate_positions(
        self, 
        num_items: int, 
        layout: str,
        content_top: float,
        content_height: float
    ) -> List[Tuple[float, float, float, float]]:
        """Calculate positions for multiple items"""
        slide_width = 13.0
        margin = 0.5
        gap = 0.3
        
        if layout == "horizontal":
            item_width = (slide_width - 2*margin - (num_items-1)*gap) / num_items
            return [
                (margin + i*(item_width + gap), content_top, item_width, content_height)
                for i in range(num_items)
            ]
        elif layout == "vertical":
            item_height = (content_height - (num_items-1)*gap) / num_items
            item_width = slide_width - 2*margin
            return [
                (margin, content_top + i*(item_height + gap), item_width, item_height)
                for i in range(num_items)
            ]
        else:  # grid
            cols = 2 if num_items <= 4 else 3
            rows = (num_items + cols - 1) // cols
            item_width = (slide_width - 2*margin - (cols-1)*gap) / cols
            item_height = (content_height - (rows-1)*gap) / rows
            
            positions = []
            for i in range(num_items):
                row = i // cols
                col = i % cols
                left = margin + col*(item_width + gap)
                top = content_top + row*(item_height + gap)
                positions.append((left, top, item_width, item_height))
            return positions
    
    def fill_template_placeholders(
        self,
        chart_mapping: Dict[str, str],
        placeholder_prefix: str = "{{CHART_"
    ) -> None:
        """
        Fill placeholder images in template with actual charts
        
        Args:
            chart_mapping: Dict mapping placeholder name to chart image path
            placeholder_prefix: Prefix used to identify chart placeholders
        """
        for slide in self.prs.slides:
            for shape in slide.shapes:
                # Check if shape name matches a placeholder
                if shape.name and placeholder_prefix in shape.name:
                    # Extract placeholder key
                    key = shape.name.replace(placeholder_prefix, "").replace("}}", "")
                    
                    if key in chart_mapping and os.path.exists(chart_mapping[key]):
                        # Get position of placeholder
                        left = shape.left
                        top = shape.top
                        width = shape.width
                        height = shape.height
                        
                        # Remove placeholder
                        sp = shape._element
                        sp.getparent().remove(sp)
                        
                        # Add chart image in same position
                        slide.shapes.add_picture(
                            chart_mapping[key],
                            left, top, width=width
                        )
    
    def generate_from_charts(
        self,
        charts: List[Dict],
        title: str = "Chart Report",
        charts_per_slide: int = 1
    ) -> None:
        """
        Generate presentation from extracted charts
        
        Args:
            charts: List of chart dicts with 'image_path', 'name', 'sheet' keys
            title: Presentation title
            charts_per_slide: Number of charts per slide
        """
        # Add title slide
        self.add_title_slide(title, f"Generated from {len(charts)} charts")
        
        # Add chart slides
        if charts_per_slide == 1:
            for chart in charts:
                self.add_chart_slide(
                    chart['image_path'],
                    title=f"{chart['sheet']} - {chart['name']}"
                )
        else:
            # Group charts
            for i in range(0, len(charts), charts_per_slide):
                group = charts[i:i+charts_per_slide]
                chart_paths = [c['image_path'] for c in group]
                group_title = f"Charts {i+1}-{i+len(group)}"
                self.add_multiple_charts_slide(chart_paths, title=group_title)
    
    def save(self, output_path: str) -> str:
        """Save the presentation"""
        output_path = os.path.abspath(output_path)
        os.makedirs(os.path.dirname(output_path), exist_ok=True)
        self.prs.save(output_path)
        return output_path


def generate_ppt_from_charts(
    charts: List[Dict],
    output_path: str,
    template_path: Optional[str] = None,
    title: str = "Chart Report",
    charts_per_slide: int = 1
) -> str:
    """
    Convenience function to generate PPT from charts
    
    Args:
        charts: List of chart dicts from ExcelChartExtractor
        output_path: Where to save the PPT
        template_path: Optional template file
        title: Presentation title
        charts_per_slide: Charts per slide
        
    Returns:
        Path to generated PPT
    """
    generator = PPTGenerator(template_path)
    generator.generate_from_charts(charts, title, charts_per_slide)
    return generator.save(output_path)

