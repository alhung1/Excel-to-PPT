"""
Netgear NRH Report Generator - Interactive Version
Allows manual input for page numbers
"""
import os
import pythoncom
import win32com.client as win32
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE


# ============================================================
# CONFIGURATION - Modify these paths as needed
# ============================================================
EXCEL_FILE = r"C:\Netgear Projects\NRH testing\RBE770v2\NRH Test Result\RBE773v2 vs RBE773 vs Eero Pro 7_3Pack Test\NRH_RBE773v2 vs RBE773 vs Eero Pro 7_auto channle_Intel BE200_20251016.xlsm"
PPT_TEMPLATE = r"C:\Netgear Projects\NRH_Report.pptx"
OUTPUT_PPT = r"C:\Netgear Projects\NRH_Report_Generated.pptx"
TEMP_DIR = r"C:\Users\alhung\excel-to-ppt\temp_charts"

# Image positioning (based on your adjusted template)
IMG_LEFT = Inches(0.423)
IMG_TOP = Inches(1.1)
IMG_WIDTH = Inches(12.0)
IMG_HEIGHT = Inches(5.6)


def list_available_charts(excel_path):
    """List all available worksheets and chart sheets"""
    pythoncom.CoInitialize()
    
    excel_app = win32.DispatchEx('Excel.Application')
    excel_app.Visible = False
    excel_app.DisplayAlerts = False
    
    workbook = excel_app.Workbooks.Open(excel_path)
    
    worksheets = []
    chartsheets = []
    
    for sheet in workbook.Worksheets:
        worksheets.append(sheet.Name)
    
    for chart in workbook.Charts:
        chartsheets.append(chart.Name)
    
    workbook.Close(SaveChanges=False)
    excel_app.Quit()
    pythoncom.CoUninitialize()
    
    return worksheets, chartsheets


def capture_item(excel_app, workbook, name, item_type, output_path):
    """Capture a worksheet or chart sheet as image"""
    try:
        if item_type == 'chartsheet':
            chart_sheet = workbook.Charts(name)
            chart_sheet.Export(output_path, "PNG")
        else:  # worksheet
            sheet = workbook.Worksheets(name)
            # Check for chart objects first
            if sheet.ChartObjects().Count > 0:
                sheet.ChartObjects(1).Chart.Export(output_path, "PNG")
            else:
                # Capture used range
                used_range = sheet.UsedRange
                used_range.CopyPicture(Appearance=1, Format=2)
                temp_chart = workbook.Charts.Add()
                temp_chart.Paste()
                temp_chart.Export(output_path, "PNG")
                excel_app.DisplayAlerts = False
                temp_chart.Delete()
        return True
    except Exception as e:
        print(f"  ✗ Error: {e}")
        return False


def main():
    print("\n" + "=" * 60)
    print("🔧 Netgear NRH Report Generator - Interactive Mode")
    print("=" * 60)
    
    # Check files
    if not os.path.exists(EXCEL_FILE):
        print(f"\n✗ Excel file not found: {EXCEL_FILE}")
        return
    if not os.path.exists(PPT_TEMPLATE):
        print(f"\n✗ PPT template not found: {PPT_TEMPLATE}")
        return
    
    print(f"\n📊 Excel: {os.path.basename(EXCEL_FILE)}")
    print(f"📄 Template: {os.path.basename(PPT_TEMPLATE)}")
    
    # List available charts
    print("\n⏳ Scanning Excel file...")
    worksheets, chartsheets = list_available_charts(EXCEL_FILE)
    
    print("\n" + "-" * 60)
    print("📋 Available items to extract:")
    print("-" * 60)
    
    all_items = []
    idx = 1
    
    print("\n[Worksheets]")
    for ws in worksheets:
        print(f"  {idx}. {ws}")
        all_items.append((ws, 'worksheet'))
        idx += 1
    
    print("\n[Chart Sheets]")
    for cs in chartsheets:
        print(f"  {idx}. {cs}")
        all_items.append((cs, 'chartsheet'))
        idx += 1
    
    # Get PPT info
    prs = Presentation(PPT_TEMPLATE)
    total_slides = len(prs.slides)
    print(f"\n📑 PPT has {total_slides} slides")
    
    # Interactive input
    print("\n" + "=" * 60)
    print("📝 Enter your selections")
    print("=" * 60)
    print("Format: <item number>,<page number>")
    print("Example: 4,8 means put item #4 on page 8")
    print("Enter 'done' when finished, 'quit' to cancel")
    print("-" * 60)
    
    selections = []
    
    while True:
        user_input = input("\n➤ Selection: ").strip().lower()
        
        if user_input == 'done':
            break
        elif user_input == 'quit':
            print("Cancelled.")
            return
        
        try:
            parts = user_input.split(',')
            item_num = int(parts[0].strip())
            page_num = int(parts[1].strip())
            
            if item_num < 1 or item_num > len(all_items):
                print(f"  ✗ Invalid item number. Must be 1-{len(all_items)}")
                continue
            if page_num < 1 or page_num > total_slides:
                print(f"  ✗ Invalid page number. Must be 1-{total_slides}")
                continue
            
            item_name, item_type = all_items[item_num - 1]
            selections.append({
                'name': item_name,
                'type': item_type,
                'page': page_num
            })
            print(f"  ✓ Added: {item_name} → Page {page_num}")
            
        except (ValueError, IndexError):
            print("  ✗ Invalid format. Use: <item>,<page>  Example: 4,8")
    
    if not selections:
        print("\n✗ No selections made. Exiting.")
        return
    
    # Confirm selections
    print("\n" + "=" * 60)
    print("📋 Your selections:")
    print("=" * 60)
    for sel in selections:
        print(f"  • {sel['name']} → Page {sel['page']}")
    
    confirm = input("\nProceed? (y/n): ").strip().lower()
    if confirm != 'y':
        print("Cancelled.")
        return
    
    # Extract charts
    print("\n" + "=" * 60)
    print("Step 1: Extracting from Excel")
    print("=" * 60)
    
    os.makedirs(TEMP_DIR, exist_ok=True)
    pythoncom.CoInitialize()
    
    excel_app = win32.DispatchEx('Excel.Application')
    excel_app.Visible = False
    excel_app.DisplayAlerts = False
    workbook = excel_app.Workbooks.Open(EXCEL_FILE)
    
    extracted = {}
    for sel in selections:
        name = sel['name']
        item_type = sel['type']
        safe_name = name.replace(' ', '_').replace('#', '_').replace('/', '_')
        output_path = os.path.join(TEMP_DIR, f"{safe_name}.png")
        
        print(f"\n→ Extracting: {name}")
        if capture_item(excel_app, workbook, name, item_type, output_path):
            extracted[name] = output_path
            print(f"  ✓ Saved: {os.path.basename(output_path)}")
    
    workbook.Close(SaveChanges=False)
    excel_app.Quit()
    pythoncom.CoUninitialize()
    
    # Insert into PPT
    print("\n" + "=" * 60)
    print("Step 2: Inserting into PowerPoint")
    print("=" * 60)
    
    prs = Presentation(PPT_TEMPLATE)
    
    for sel in selections:
        name = sel['name']
        page = sel['page']
        slide_idx = page - 1
        
        if name not in extracted:
            print(f"\n✗ {name} - no image")
            continue
        
        print(f"\n→ {name} → Page {page}")
        
        slide = prs.slides[slide_idx]
        image_path = extracted[name]
        
        # Add picture with exact positioning
        pic = slide.shapes.add_picture(
            image_path,
            IMG_LEFT,
            IMG_TOP,
            width=IMG_WIDTH,
            height=IMG_HEIGHT
        )
        print(f"  ✓ Inserted at ({IMG_LEFT.inches:.2f}\", {IMG_TOP.inches:.2f}\") size {IMG_WIDTH.inches:.1f}\"x{IMG_HEIGHT.inches:.1f}\"")
    
    # Save
    prs.save(OUTPUT_PPT)
    print(f"\n✓ Saved: {OUTPUT_PPT}")
    
    print("\n" + "=" * 60)
    print(f"🎉 Done! Output: {OUTPUT_PPT}")
    print("=" * 60)


if __name__ == "__main__":
    main()

