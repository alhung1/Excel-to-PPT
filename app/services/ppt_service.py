"""
PowerPoint generation service.

Handles both image-mode and embedded-mode chart insertion.
"""
import os
import time
import shutil
from pathlib import Path
from typing import Dict, List, Optional

from pptx import Presentation
from pptx.util import Inches

from app.config import (
    logger,
    MESH_BACKHAUL_LAYOUT,
    MESH_FRONTHAUL_LAYOUT,
    COM_CLIPBOARD_DELAY,
)
from app.models.schemas import ChartMapping, GenerateRequest
from app.services.excel_service import ExcelCOM, PowerPointCOM, capture_item
from app.utils.clipboard import clear_clipboard


# ---------------------------------------------------------------------------
# Slide helpers
# ---------------------------------------------------------------------------
def get_slide_title(slide) -> str:
    """Read the title text from a python-pptx slide."""
    try:
        if slide.shapes.title:
            return (slide.shapes.title.text or "").strip()
    except Exception:
        pass
    return ""


def get_ppt_info(ppt_path: str) -> dict:
    """Return slide metadata for a PPT file."""
    prs = Presentation(ppt_path)
    slides = []
    for idx, slide in enumerate(prs.slides):
        slides.append({"page": idx + 1, "title": get_slide_title(slide)})
    return {
        "total_slides": len(prs.slides),
        "slides": slides,
        "width": prs.slide_width.inches,
        "height": prs.slide_height.inches,
    }


def get_ppt_slide_titles(ppt_path: str) -> Dict[int, str]:
    """Return a {page_number: title} mapping."""
    prs = Presentation(ppt_path)
    return {
        idx + 1: get_slide_title(slide) for idx, slide in enumerate(prs.slides)
    }


# ---------------------------------------------------------------------------
# Layout helpers
# ---------------------------------------------------------------------------
def get_mesh_layout_type(title: str) -> Optional[str]:
    """Detect mesh layout type from slide title."""
    if not title:
        return None
    t = title.lower()
    if "mesh" not in t:
        return None
    if "backhaul" in t:
        return "backhaul"
    if "fronthaul" in t:
        return "fronthaul"
    return None


def is_mesh_slide_title(title: str) -> bool:
    return get_mesh_layout_type(title) is not None


def get_effective_layout(request: GenerateRequest, slide_title: str) -> Dict[str, float]:
    """Return the image placement box for a slide."""
    mesh_type = get_mesh_layout_type(slide_title)
    if mesh_type == "backhaul":
        return dict(MESH_BACKHAUL_LAYOUT)
    if mesh_type == "fronthaul":
        return dict(MESH_FRONTHAUL_LAYOUT)
    return {
        "left": request.img_left,
        "top": request.img_top,
        "width": request.img_width,
        "height": request.img_height,
    }


# ---------------------------------------------------------------------------
# Image-mode processing
# ---------------------------------------------------------------------------
def process_image_mappings(
    mappings: List[ChartMapping],
    prs: Presentation,
    request: GenerateRequest,
    job_dir: Path,
    slide_titles: Dict[int, str],
    uploaded_files: dict,
) -> List[dict]:
    """Insert charts as static PNG images into a python-pptx Presentation."""
    results: List[dict] = []

    # Group by Excel file
    excel_files: Dict[str, dict] = {}
    for m in mappings:
        if m.excel_id not in excel_files:
            excel_files[m.excel_id] = {
                "path": uploaded_files[m.excel_id]["path"],
                "filename": uploaded_files[m.excel_id]["filename"],
                "mappings": [],
            }
        excel_files[m.excel_id]["mappings"].append(m)

    # Extract images
    extracted: Dict[str, str] = {}
    with ExcelCOM() as (excel_app, _):
        for excel_id, info in excel_files.items():
            logger.info("[Image Mode] Opening: %s", info["filename"])
            workbook = excel_app.Workbooks.Open(info["path"])

            for mapping in info["mappings"]:
                key = f"{excel_id}|{mapping.name}"
                safe_name = _safe_filename(f"{excel_id}_{mapping.name}")
                out_path = str(job_dir / f"{safe_name}.png")

                logger.info("  Capturing: %s (type: %s)", mapping.name, mapping.type)
                if capture_item(excel_app, workbook, mapping.name, mapping.type, out_path):
                    extracted[key] = out_path
                    logger.info("  [OK] Extracted: %s (%d bytes)", mapping.name, os.path.getsize(out_path))
                else:
                    logger.warning("  [FAIL] Failed to extract: %s", mapping.name)

            workbook.Close(SaveChanges=False)

    # Insert into PPT
    for mapping in mappings:
        key = f"{mapping.excel_id}|{mapping.name}"
        excel_filename = uploaded_files[mapping.excel_id]["filename"]

        if key not in extracted:
            results.append({"name": mapping.name, "excel": excel_filename, "status": "failed", "reason": "擷取失敗"})
            continue

        slide_idx = mapping.page - 1
        if slide_idx >= len(prs.slides):
            results.append({"name": mapping.name, "excel": excel_filename, "status": "failed", "reason": f"第 {mapping.page} 頁不存在"})
            continue

        image_path = extracted[key]
        if not os.path.exists(image_path):
            results.append({"name": mapping.name, "excel": excel_filename, "status": "failed", "reason": "圖片檔案不存在"})
            continue

        image_size = os.path.getsize(image_path)
        if image_size < 500:
            results.append({"name": mapping.name, "excel": excel_filename, "status": "failed", "reason": f"圖片檔案可能損壞 (大小: {image_size} bytes)"})
            continue

        try:
            slide = prs.slides[slide_idx]
            slide_title = slide_titles.get(mapping.page, "")
            layout = get_effective_layout(request, slide_title)
            slide.shapes.add_picture(
                image_path,
                Inches(layout["left"]),
                Inches(layout["top"]),
                width=Inches(layout["width"]),
                height=Inches(layout["height"]),
            )
            results.append({
                "name": mapping.name,
                "excel": excel_filename,
                "status": "success",
                "page": mapping.page,
                "mode": "image",
                "mesh_layout": is_mesh_slide_title(slide_title),
            })
        except Exception as e:
            logger.error("Error adding image to slide: %s", e)
            results.append({"name": mapping.name, "excel": excel_filename, "status": "failed", "reason": f"無法插入圖片: {e}"})

    return results


# ---------------------------------------------------------------------------
# Embedded-mode processing
# ---------------------------------------------------------------------------
def process_embedded_mappings(
    mappings: List[ChartMapping],
    ppt_path: str,
    request: GenerateRequest,
    slide_titles: Dict[int, str],
    uploaded_files: dict,
) -> List[dict]:
    """Insert charts as editable objects using COM copy-paste."""
    results: List[dict] = []

    # Group by Excel file
    excel_files: Dict[str, dict] = {}
    for m in mappings:
        if m.excel_id not in excel_files:
            excel_files[m.excel_id] = {
                "path": uploaded_files[m.excel_id]["path"],
                "filename": uploaded_files[m.excel_id]["filename"],
                "mappings": [],
            }
        excel_files[m.excel_id]["mappings"].append(m)

    # We need both Excel and PowerPoint COM, both visible
    import pythoncom
    import win32com.client as win32

    pythoncom.CoInitialize()
    excel_app = None
    ppt_app = None
    presentation = None

    try:
        logger.info("[Embedded Mode] Starting Excel...")
        excel_app = win32.DispatchEx("Excel.Application")
        excel_app.Visible = True
        excel_app.DisplayAlerts = False

        logger.info("[Embedded Mode] Starting PowerPoint...")
        ppt_app = win32.DispatchEx("PowerPoint.Application")
        ppt_app.Visible = True

        presentation = ppt_app.Presentations.Open(ppt_path)
        time.sleep(0.5)

        for excel_id, info in excel_files.items():
            logger.info("[Embedded Mode] Opening: %s", info["filename"])
            workbook = excel_app.Workbooks.Open(os.path.abspath(info["path"]))
            time.sleep(0.3)

            for mapping in info["mappings"]:
                excel_filename = info["filename"]
                logger.info("  [Embedded] Processing: %s -> Page %d", mapping.name, mapping.page)

                if mapping.page > presentation.Slides.Count:
                    results.append({"name": mapping.name, "excel": excel_filename, "status": "failed", "reason": f"第 {mapping.page} 頁不存在"})
                    continue

                try:
                    if mapping.type == "chartsheet":
                        chart_sheet = workbook.Charts(mapping.name)
                        chart_sheet.Activate()
                        time.sleep(0.2)
                        clear_clipboard()
                        time.sleep(0.1)
                        chart_sheet.ChartArea.Copy()
                    else:
                        sheet = workbook.Worksheets(mapping.name)
                        sheet.Activate()
                        time.sleep(0.2)

                        chart_count = 0
                        try:
                            chart_count = sheet.ChartObjects().Count
                        except Exception:
                            pass

                        if chart_count > 0:
                            chart_obj = sheet.ChartObjects(1)
                            chart_obj.Select()
                            time.sleep(0.2)
                            clear_clipboard()
                            time.sleep(0.1)
                            chart_obj.Chart.ChartArea.Copy()
                        else:
                            logger.info("    [Embedded] No chart found, skipping")
                            results.append({"name": mapping.name, "excel": excel_filename, "status": "failed", "reason": "工作表中沒有圖表"})
                            continue

                    time.sleep(0.5)

                    slide = presentation.Slides(mapping.page)
                    try:
                        shape = slide.Shapes.Paste()
                        time.sleep(0.3)
                        if hasattr(shape, "Item"):
                            shape = shape.Item(1)

                        slide_title = slide_titles.get(mapping.page, "")
                        layout = get_effective_layout(request, slide_title)
                        shape.Left = layout["left"] * 72
                        shape.Top = layout["top"] * 72
                        shape.Width = layout["width"] * 72
                        shape.Height = layout["height"] * 72

                        shape.Line.Visible = -1
                        shape.Line.ForeColor.RGB = 0
                        shape.Line.Weight = 0.75

                        results.append({
                            "name": mapping.name,
                            "excel": excel_filename,
                            "status": "success",
                            "page": mapping.page,
                            "mode": "embedded",
                            "mesh_layout": is_mesh_slide_title(slide_title),
                        })
                        logger.info("    [OK] Embedded chart pasted successfully")

                    except Exception as paste_error:
                        logger.warning("    [FAIL] Paste failed: %s", paste_error)
                        results.append({"name": mapping.name, "excel": excel_filename, "status": "failed", "reason": f"貼上失敗: {paste_error}"})

                except Exception as e:
                    logger.error("    [ERROR] Processing %s: %s", mapping.name, e)
                    results.append({"name": mapping.name, "excel": excel_filename, "status": "failed", "reason": str(e)})

            workbook.Close(SaveChanges=False)

        presentation.Save()
        presentation.Close()
        presentation = None

    except Exception as e:
        logger.error("Embedded mode error: %s", e, exc_info=True)
        raise
    finally:
        if presentation:
            try:
                presentation.Close()
            except Exception:
                pass
        if excel_app:
            try:
                excel_app.Quit()
            except Exception:
                pass
        if ppt_app:
            try:
                ppt_app.Quit()
            except Exception:
                pass
        pythoncom.CoUninitialize()

    return results


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------
def _safe_filename(name: str) -> str:
    """Sanitise a string for use as a file name."""
    for char in '<>:"/\\|?*# ':
        name = name.replace(char, "_")
    return name
