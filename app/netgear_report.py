"""
Netgear NRH Report Generator
Extracts specific sheets/charts from Excel and places them in PPT template
"""
import os
import pythoncom
import win32com.client as win32
from pptx import Presentation
from pptx.util import Inches, Pt
import shutil
from pathlib import Path


# Configuration
EXCEL_FILE = r"C:\Netgear Projects\NRH testing\RBE770v2\NRH Test Result\RBE773v2 vs RBE773 vs Eero Pro 7_3Pack Test\NRH_RBE773v2 vs RBE773 vs Eero Pro 7_auto channle_Intel BE200_20251016.xlsm"
PPT_TEMPLATE = r"C:\Netgear Projects\NRH_Report.pptx"
OUTPUT_PPT = r"C:\Netgear Projects\NRH_Report_Generated.pptx"

# Mapping: chart/sheet name -> (PPT slide index 0-based, type)
# type: 'worksheet' for regular sheets, 'chartsheet' for standalone chart sheets
ITEMS_TO_EXTRACT = {
    "Metric DUT vs REF#1": {"slide": 7, "type": "worksheet"},   # Page 8 - Matrix table
    "BI": {"slide": 8, "type": "chartsheet"},                    # Page 9 - UL+DL
    "DL": {"slide": 9, "type": "chartsheet"},                    # Page 10
    "UL": {"slide": 10, "type": "chartsheet"},                   # Page 11
}

# Temp directory for chart images
TEMP_DIR = r"C:\Users\alhung\excel-to-ppt\temp_charts"


def capture_worksheet(excel_app, workbook, sheet_name, output_path):
    """Capture a regular worksheet's used range as image"""
    try:
        sheet = workbook.Worksheets(sheet_name)
        print(f"  Found worksheet: {sheet_name}")
        
        # Check for chart objects first
        chart_count = sheet.ChartObjects().Count
        if chart_count > 0:
            # Export the first chart
            chart_obj = sheet.ChartObjects(1)
            chart_obj.Chart.Export(output_path, "PNG")
            print(f"  Exported chart object to: {output_path}")
            return True
        
        # No chart objects, capture the used range
        print(f"  Capturing used range as image...")
        used_range = sheet.UsedRange
        
        # Copy as picture
        used_range.CopyPicture(Appearance=1, Format=2)  # xlScreen=1, xlBitmap=2
        
        # Create a temporary chart sheet to paste and export
        temp_chart_sheet = workbook.Charts.Add()
        temp_chart_sheet.Paste()
        temp_chart_sheet.Export(output_path, "PNG")
        
        # Delete temp chart sheet
        excel_app.DisplayAlerts = False
        temp_chart_sheet.Delete()
        excel_app.DisplayAlerts = True
        
        print(f"  Exported to: {output_path}")
        return True
        
    except Exception as e:
        print(f"  Error: {e}")
        return False


def capture_chartsheet(workbook, chart_name, output_path):
    """Capture a standalone chart sheet as image"""
    try:
        # Access chart sheets via workbook.Charts collection
        chart_sheet = workbook.Charts(chart_name)
        print(f"  Found chart sheet: {chart_name}")
        
        # Export directly
        chart_sheet.Export(output_path, "PNG")
        print(f"  Exported to: {output_path}")
        return True
        
    except Exception as e:
        print(f"  Error: {e}")
        return False


def extract_charts_from_excel():
    """Extract charts from the specified Excel sheets"""
    print("=" * 60)
    print("Step 1: Extracting charts from Excel")
    print("=" * 60)
    
    # Create temp directory
    os.makedirs(TEMP_DIR, exist_ok=True)
    
    # Initialize COM
    pythoncom.CoInitialize()
    
    excel_app = None
    workbook = None
    extracted_images = {}
    
    try:
        # Start Excel
        print(f"\nOpening Excel...")
        excel_app = win32.DispatchEx('Excel.Application')
        excel_app.Visible = False
        excel_app.DisplayAlerts = False
        
        # Open workbook
        print(f"Opening: {os.path.basename(EXCEL_FILE)}")
        workbook = excel_app.Workbooks.Open(EXCEL_FILE)
        
        # Extract each specified item
        print(f"\nExtracting items:")
        for name, config in ITEMS_TO_EXTRACT.items():
            print(f"\nProcessing: {name} (type: {config['type']})")
            
            # Safe filename
            safe_name = name.replace(' ', '_').replace('#', '_')
            output_path = os.path.join(TEMP_DIR, f"{safe_name}.png")
            
            if config['type'] == 'worksheet':
                success = capture_worksheet(excel_app, workbook, name, output_path)
            else:  # chartsheet
                success = capture_chartsheet(workbook, name, output_path)
            
            if success:
                extracted_images[name] = output_path
        
        print(f"\n✓ Extracted {len(extracted_images)} images")
        return extracted_images
        
    except Exception as e:
        print(f"Error: {e}")
        import traceback
        traceback.print_exc()
        return {}
        
    finally:
        if workbook:
            workbook.Close(SaveChanges=False)
        if excel_app:
            excel_app.Quit()
        pythoncom.CoUninitialize()


def insert_images_to_ppt(extracted_images):
    """Insert extracted images into the PPT template"""
    print("\n" + "=" * 60)
    print("Step 2: Inserting images into PowerPoint")
    print("=" * 60)
    
    # Open template
    print(f"\nOpening template: {os.path.basename(PPT_TEMPLATE)}")
    prs = Presentation(PPT_TEMPLATE)
    
    # Get slide dimensions
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    
    print(f"Slide size: {slide_width.inches:.2f}\" x {slide_height.inches:.2f}\"")
    print(f"Total slides: {len(prs.slides)}")
    
    # Define content area (leaving space for title)
    margin_left = Inches(0.3)
    margin_top = Inches(1.3)       # Space for title
    margin_right = Inches(0.3)
    margin_bottom = Inches(0.3)
    
    content_width = slide_width - margin_left - margin_right
    content_height = slide_height - margin_top - margin_bottom
    
    # Insert images
    for name, config in ITEMS_TO_EXTRACT.items():
        slide_idx = config['slide']
        
        if name not in extracted_images:
            print(f"\n✗ Skipping {name} - no image")
            continue
            
        image_path = extracted_images[name]
        
        if slide_idx >= len(prs.slides):
            print(f"\n✗ Skipping {name} - slide {slide_idx + 1} doesn't exist")
            continue
        
        print(f"\n→ Inserting {name} into slide {slide_idx + 1}")
        
        slide = prs.slides[slide_idx]
        
        try:
            # Add picture
            pic = slide.shapes.add_picture(
                image_path,
                margin_left,
                margin_top,
                width=content_width
            )
            
            # If too tall, scale to fit height instead
            if pic.height > content_height:
                scale = float(content_height) / float(pic.height)
                pic.width = int(pic.width * scale)
                pic.height = content_height
                # Re-center horizontally
                pic.left = int((slide_width - pic.width) / 2)
            else:
                # Center horizontally
                pic.left = int((slide_width - pic.width) / 2)
            
            print(f"  ✓ Size: {pic.width.inches:.2f}\" x {pic.height.inches:.2f}\"")
            
        except Exception as e:
            print(f"  ✗ Error: {e}")
    
    # Save output
    print(f"\nSaving: {OUTPUT_PPT}")
    prs.save(OUTPUT_PPT)
    print("✓ Done!")
    
    return OUTPUT_PPT


def main():
    """Main function"""
    print("\n" + "=" * 60)
    print("Netgear NRH Report Generator")
    print("=" * 60)
    print(f"\nExcel: {os.path.basename(EXCEL_FILE)}")
    print(f"Template: {os.path.basename(PPT_TEMPLATE)}")
    print(f"Output: {os.path.basename(OUTPUT_PPT)}")
    
    # Check files exist
    if not os.path.exists(EXCEL_FILE):
        print(f"\n✗ Excel file not found!")
        return
    if not os.path.exists(PPT_TEMPLATE):
        print(f"\n✗ PPT template not found!")
        return
    
    # Extract charts
    extracted_images = extract_charts_from_excel()
    
    if not extracted_images:
        print("\n✗ No images extracted. Exiting.")
        return
    
    # Insert into PPT
    output_path = insert_images_to_ppt(extracted_images)
    
    print("\n" + "=" * 60)
    print(f"✓ Report generated: {output_path}")
    print("=" * 60)


if __name__ == "__main__":
    main()
