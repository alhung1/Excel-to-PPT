"""
Analyze PPT to get image positions and sizes
"""
from pptx import Presentation
from pptx.util import Inches, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE

PPT_FILE = r"C:\Netgear Projects\NRH_Report_Generatedupdate.pptx"

prs = Presentation(PPT_FILE)

print("=" * 70)
print(f"Analyzing: {PPT_FILE}")
print(f"Slide size: {prs.slide_width.inches:.2f}\" x {prs.slide_height.inches:.2f}\"")
print("=" * 70)

# Check slides 8-11 (index 7-10)
for slide_idx in range(7, min(12, len(prs.slides))):
    slide = prs.slides[slide_idx]
    print(f"\n{'='*70}")
    print(f"📄 Slide {slide_idx + 1}:")
    
    # Get title
    if slide.shapes.title:
        print(f"   Title: {slide.shapes.title.text}")
    
    # Find pictures (skip small decorative elements)
    pic_count = 0
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            # Skip small images (likely decorative)
            if shape.width.inches < 2 or shape.height.inches < 2:
                continue
            
            pic_count += 1
            print(f"\n   🖼️ Picture {pic_count}:")
            print(f"      Position: left={shape.left.inches:.3f}\", top={shape.top.inches:.3f}\"")
            print(f"      Size:     width={shape.width.inches:.3f}\", height={shape.height.inches:.3f}\"")
    
    if pic_count == 0:
        print(f"   (No large pictures found)")
