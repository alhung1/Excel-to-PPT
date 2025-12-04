"""
Excel to PowerPoint Web Application
Full-featured GUI - File Upload Support
"""
import os
import uuid
import shutil
from datetime import datetime
from pathlib import Path
from typing import List, Optional, Dict
import pythoncom
import win32com.client as win32

from fastapi import FastAPI, UploadFile, File, Form, HTTPException
from fastapi.responses import FileResponse, HTMLResponse
from fastapi.staticfiles import StaticFiles
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
from pptx import Presentation
from pptx.util import Inches

# Create app
app = FastAPI(
    title="Excel to PowerPoint Generator",
    description="Upload Excel and PPT files to generate reports",
    version="4.0.0"
)

# CORS middleware
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Directories
BASE_DIR = Path(__file__).parent.parent
UPLOAD_DIR = BASE_DIR / "uploads"
OUTPUT_DIR = BASE_DIR / "outputs"
STATIC_DIR = BASE_DIR / "static"

# Ensure directories exist
for d in [UPLOAD_DIR, OUTPUT_DIR, STATIC_DIR]:
    d.mkdir(exist_ok=True)

# Mount static files
app.mount("/static", StaticFiles(directory=str(STATIC_DIR)), name="static")

# Store uploaded files info (in-memory for simplicity)
uploaded_files = {}


# ============================================================
# Pydantic Models
# ============================================================
class ChartMapping(BaseModel):
    excel_id: str         # Uploaded Excel file ID
    name: str             # Sheet/Chart name
    page: int             # Target page number
    type: str             # "worksheet" or "chartsheet"


class GenerateRequest(BaseModel):
    template_id: str
    output_name: str
    mappings: List[ChartMapping]
    img_left: float = 0.423
    img_top: float = 1.1
    img_width: float = 12.0
    img_height: float = 5.6


# ============================================================
# Excel Functions
# ============================================================
def get_excel_info(excel_path: str) -> dict:
    """Get worksheets and chart sheets from Excel file"""
    pythoncom.CoInitialize()
    
    try:
        excel_app = win32.DispatchEx('Excel.Application')
        excel_app.Visible = False
        excel_app.DisplayAlerts = False
        
        workbook = excel_app.Workbooks.Open(excel_path)
        
        worksheets = []
        for sheet in workbook.Worksheets:
            chart_count = 0
            try:
                chart_count = sheet.ChartObjects().Count
            except:
                pass
            worksheets.append({
                "name": sheet.Name,
                "type": "worksheet",
                "has_charts": chart_count > 0,
                "chart_count": chart_count
            })
        
        chartsheets = []
        try:
            for chart in workbook.Charts:
                chartsheets.append({
                    "name": chart.Name,
                    "type": "chartsheet"
                })
        except:
            pass
        
        workbook.Close(SaveChanges=False)
        excel_app.Quit()
        
        return {
            "worksheets": worksheets,
            "chartsheets": chartsheets
        }
    except Exception as e:
        raise e
    finally:
        pythoncom.CoUninitialize()


def capture_item(excel_app, workbook, name: str, item_type: str, output_path: str) -> bool:
    """Capture a worksheet or chart sheet as image"""
    try:
        if item_type == 'chartsheet':
            chart_sheet = workbook.Charts(name)
            chart_sheet.Export(output_path, "PNG")
        else:
            sheet = workbook.Worksheets(name)
            if sheet.ChartObjects().Count > 0:
                sheet.ChartObjects(1).Chart.Export(output_path, "PNG")
            else:
                used_range = sheet.UsedRange
                used_range.CopyPicture(Appearance=1, Format=2)
                temp_chart = workbook.Charts.Add()
                temp_chart.Paste()
                temp_chart.Export(output_path, "PNG")
                excel_app.DisplayAlerts = False
                temp_chart.Delete()
        return True
    except Exception as e:
        print(f"Error capturing {name}: {e}")
        return False


def get_ppt_info(ppt_path: str) -> dict:
    """Get PPT slide info"""
    prs = Presentation(ppt_path)
    slides = []
    for idx, slide in enumerate(prs.slides):
        title = ""
        if slide.shapes.title:
            title = slide.shapes.title.text
        slides.append({
            "page": idx + 1,
            "title": title
        })
    return {
        "total_slides": len(prs.slides),
        "slides": slides,
        "width": prs.slide_width.inches,
        "height": prs.slide_height.inches
    }


# ============================================================
# API Endpoints
# ============================================================
@app.get("/", response_class=HTMLResponse)
async def home():
    """Serve the main HTML page"""
    html_path = STATIC_DIR / "index.html"
    if html_path.exists():
        return html_path.read_text(encoding='utf-8')
    return "<html><body><h1>Please add index.html</h1></body></html>"


@app.post("/api/upload-excel")
async def upload_excel(file: UploadFile = File(...)):
    """Upload Excel file and get sheet info"""
    if not file.filename.endswith(('.xlsx', '.xlsm', '.xls')):
        raise HTTPException(400, "必須是 Excel 檔案 (.xlsx, .xlsm, .xls)")
    
    file_id = uuid.uuid4().hex[:8]
    file_path = UPLOAD_DIR / f"{file_id}_{file.filename}"
    
    # Save uploaded file
    with open(file_path, "wb") as f:
        content = await file.read()
        f.write(content)
    
    try:
        info = get_excel_info(str(file_path))
        
        # Store file info
        uploaded_files[file_id] = {
            "type": "excel",
            "path": str(file_path),
            "filename": file.filename
        }
        
        return {
            "status": "success",
            "file_id": file_id,
            "filename": file.filename,
            "worksheets": info["worksheets"],
            "chartsheets": info["chartsheets"]
        }
    except Exception as e:
        # Clean up on error
        if file_path.exists():
            file_path.unlink()
        raise HTTPException(500, f"讀取 Excel 失敗: {str(e)}")


@app.post("/api/upload-ppt")
async def upload_ppt(file: UploadFile = File(...)):
    """Upload PPT template and get slide info"""
    if not file.filename.endswith(('.pptx', '.ppt')):
        raise HTTPException(400, "必須是 PowerPoint 檔案 (.pptx, .ppt)")
    
    file_id = uuid.uuid4().hex[:8]
    file_path = UPLOAD_DIR / f"{file_id}_{file.filename}"
    
    # Save uploaded file
    with open(file_path, "wb") as f:
        content = await file.read()
        f.write(content)
    
    try:
        info = get_ppt_info(str(file_path))
        
        # Store file info
        uploaded_files[file_id] = {
            "type": "ppt",
            "path": str(file_path),
            "filename": file.filename
        }
        
        return {
            "status": "success",
            "file_id": file_id,
            "filename": file.filename,
            **info
        }
    except Exception as e:
        # Clean up on error
        if file_path.exists():
            file_path.unlink()
        raise HTTPException(500, f"讀取 PPT 失敗: {str(e)}")


@app.delete("/api/remove-file/{file_id}")
async def remove_file(file_id: str):
    """Remove an uploaded file"""
    if file_id in uploaded_files:
        file_info = uploaded_files[file_id]
        file_path = Path(file_info["path"])
        if file_path.exists():
            file_path.unlink()
        del uploaded_files[file_id]
        return {"status": "success"}
    raise HTTPException(404, "檔案不存在")


@app.post("/api/generate")
async def generate_ppt(request: GenerateRequest):
    """Generate PowerPoint with chart mappings"""
    
    # Get template path
    if request.template_id not in uploaded_files:
        raise HTTPException(404, "PPT 模板不存在，請重新上傳")
    
    template_path = uploaded_files[request.template_id]["path"]
    
    # Group mappings by Excel file
    excel_files = {}
    for m in request.mappings:
        if m.excel_id not in uploaded_files:
            raise HTTPException(404, f"Excel 檔案不存在: {m.excel_id}")
        
        if m.excel_id not in excel_files:
            excel_files[m.excel_id] = {
                "path": uploaded_files[m.excel_id]["path"],
                "filename": uploaded_files[m.excel_id]["filename"],
                "mappings": []
            }
        excel_files[m.excel_id]["mappings"].append(m)
    
    job_id = uuid.uuid4().hex[:8]
    job_dir = OUTPUT_DIR / job_id
    job_dir.mkdir(exist_ok=True)
    
    try:
        extracted = {}
        
        # Extract from each Excel file
        pythoncom.CoInitialize()
        
        excel_app = win32.DispatchEx('Excel.Application')
        excel_app.Visible = False
        excel_app.DisplayAlerts = False
        
        for excel_id, excel_info in excel_files.items():
            print(f"Opening: {excel_info['filename']}")
            workbook = excel_app.Workbooks.Open(excel_info['path'])
            
            for mapping in excel_info['mappings']:
                key = f"{excel_id}|{mapping.name}"
                safe_name = f"{excel_id}_{mapping.name}".replace(' ', '_').replace('#', '_').replace('/', '_')
                output_path = str(job_dir / f"{safe_name}.png")
                
                if capture_item(excel_app, workbook, mapping.name, mapping.type, output_path):
                    extracted[key] = output_path
                    print(f"  Extracted: {mapping.name}")
            
            workbook.Close(SaveChanges=False)
        
        excel_app.Quit()
        pythoncom.CoUninitialize()
        
        # Generate PPT
        prs = Presentation(template_path)
        
        results = []
        for mapping in request.mappings:
            key = f"{mapping.excel_id}|{mapping.name}"
            excel_filename = uploaded_files[mapping.excel_id]["filename"]
            
            if key not in extracted:
                results.append({
                    "name": mapping.name,
                    "excel": excel_filename,
                    "status": "failed",
                    "reason": "擷取失敗"
                })
                continue
            
            slide_idx = mapping.page - 1
            if slide_idx >= len(prs.slides):
                results.append({
                    "name": mapping.name,
                    "excel": excel_filename,
                    "status": "failed",
                    "reason": f"第 {mapping.page} 頁不存在"
                })
                continue
            
            slide = prs.slides[slide_idx]
            pic = slide.shapes.add_picture(
                extracted[key],
                Inches(request.img_left),
                Inches(request.img_top),
                width=Inches(request.img_width),
                height=Inches(request.img_height)
            )
            results.append({
                "name": mapping.name,
                "excel": excel_filename,
                "status": "success",
                "page": mapping.page
            })
        
        # Save output
        output_filename = f"{request.output_name}.pptx" if not request.output_name.endswith('.pptx') else request.output_name
        output_path = job_dir / output_filename
        prs.save(str(output_path))
        
        return {
            "status": "success",
            "job_id": job_id,
            "download_url": f"/api/download/{job_id}/{output_filename}",
            "results": results,
            "output_file": str(output_path)
        }
        
    except Exception as e:
        import traceback
        traceback.print_exc()
        raise HTTPException(500, f"產生 PPT 失敗: {str(e)}")


@app.get("/api/download/{job_id}/{filename}")
async def download_file(job_id: str, filename: str):
    """Download generated file"""
    file_path = OUTPUT_DIR / job_id / filename
    
    if not file_path.exists():
        raise HTTPException(404, "檔案不存在或已過期")
    
    return FileResponse(
        file_path,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        filename=filename
    )


if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)
