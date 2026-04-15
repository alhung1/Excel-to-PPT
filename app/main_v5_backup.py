"""
Excel to PowerPoint Web Application
Full-featured GUI - File Upload Support
"""
import os
import uuid
import shutil
from datetime import datetime
from pathlib import Path
from typing import List, Optional, Dict
import pythoncom
import win32com.client as win32

from fastapi import FastAPI, UploadFile, File, Form, HTTPException
from fastapi.responses import FileResponse, HTMLResponse
from fastapi.staticfiles import StaticFiles
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
from pptx import Presentation
from pptx.util import Inches

# Create app
app = FastAPI(
    title="Excel to PowerPoint Generator",
    description="Upload Excel and PPT files to generate reports",
    version="5.1.0"
)

# CORS middleware
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Directories
BASE_DIR = Path(__file__).parent.parent
UPLOAD_DIR = BASE_DIR / "uploads"
OUTPUT_DIR = BASE_DIR / "outputs"
STATIC_DIR = BASE_DIR / "static"

# Ensure directories exist
for d in [UPLOAD_DIR, OUTPUT_DIR, STATIC_DIR]:
    d.mkdir(exist_ok=True)

# Mount static files
app.mount("/static", StaticFiles(directory=str(STATIC_DIR)), name="static")

# Store uploaded files info (in-memory for simplicity)
uploaded_files = {}

# Fixed main-chart layouts derived from `templates/Mesh.pptx`.
MESH_BACKHAUL_LAYOUT = {
    "left": 0.423,
    "top": 1.267,
    "width": 12.0,
    "height": 4.294,
}

MESH_FRONTHAUL_LAYOUT = {
    "left": 0.208,
    "top": 1.267,
    "width": 12.592,
    "height": 4.344,
}


# ============================================================
# Pydantic Models
# ============================================================
class ChartMapping(BaseModel):
    excel_id: str         # Uploaded Excel file ID
    name: str             # Sheet/Chart name
    page: int             # Target page number
    type: str             # "worksheet" or "chartsheet"
    chart_mode: str = "image"  # "image" or "embedded" - per-mapping mode


class GenerateRequest(BaseModel):
    template_id: str
    output_name: str
    mappings: List[ChartMapping]
    chart_mode: str = "image"  # "image" or "embedded"
    img_left: float = 0.423
    img_top: float = 1.4
    img_width: float = 12.0
    img_height: float = 5.6


# ============================================================
# Excel Functions
# ============================================================
def get_excel_info(excel_path: str) -> dict:
    """Get worksheets and chart sheets from Excel file"""
    pythoncom.CoInitialize()
    
    try:
        excel_app = win32.DispatchEx('Excel.Application')
        excel_app.Visible = False
        excel_app.DisplayAlerts = False
        
        workbook = excel_app.Workbooks.Open(excel_path)
        
        worksheets = []
        for sheet in workbook.Worksheets:
            chart_count = 0
            try:
                chart_count = sheet.ChartObjects().Count
            except:
                pass
            worksheets.append({
                "name": sheet.Name,
                "type": "worksheet",
                "has_charts": chart_count > 0,
                "chart_count": chart_count
            })
        
        chartsheets = []
        try:
            for chart in workbook.Charts:
                chartsheets.append({
                    "name": chart.Name,
                    "type": "chartsheet"
                })
        except:
            pass
        
        workbook.Close(SaveChanges=False)
        excel_app.Quit()
        
        return {
            "worksheets": worksheets,
            "chartsheets": chartsheets
        }
    except Exception as e:
        raise e
    finally:
        pythoncom.CoUninitialize()


def capture_item(excel_app, workbook, name: str, item_type: str, output_path: str, max_retries: int = 3) -> bool:
    """Capture a worksheet or chart sheet as image
    
    Args:
        excel_app: Excel application instance
        workbook: Workbook object
        name: Sheet or chart name
        item_type: 'chartsheet' or 'worksheet'
        output_path: Path to save the PNG image
        max_retries: Maximum retry attempts for clipboard operations
    
    Returns:
        True if capture succeeded, False otherwise
    """
    import time
    
    def validate_image(path: str, min_size: int = 1000) -> bool:
        """Check if image file exists, has reasonable size, and has actual content"""
        if not os.path.exists(path):
            print(f"    [Validate] File does not exist: {path}")
            return False
        
        size = os.path.getsize(path)
        if size < min_size:
            print(f"    [Validate] File too small: {size} bytes (min: {min_size})")
            return False
        
        # Check image content using PIL
        try:
            from PIL import Image
            import statistics
            
            img = Image.open(path)
            # Convert to grayscale and get pixel data
            gray = img.convert('L')
            pixels = list(gray.getdata())
            
            # Check if image is mostly blank (very low variance)
            unique_colors = len(set(pixels))
            if unique_colors < 10:
                print(f"    [Validate] Image appears blank: only {unique_colors} unique colors")
                return False
            
            # Check standard deviation - if too low, image is likely blank/single color
            try:
                stdev = statistics.stdev(pixels)
                if stdev < 5:
                    print(f"    [Validate] Image has very low variance (stdev={stdev:.2f}), likely blank")
                    return False
            except statistics.StatisticsError:
                pass  # Not enough data points
            
            print(f"    [Validate] Image OK: {size} bytes, {unique_colors} colors")
            img.close()
            
        except Exception as e:
            print(f"    [Validate] PIL check failed ({e}), relying on file size only")
        
        return True
    
    def clear_clipboard():
        """Clear clipboard to avoid stale data issues"""
        try:
            import win32clipboard
            win32clipboard.OpenClipboard()
            win32clipboard.EmptyClipboard()
            win32clipboard.CloseClipboard()
        except:
            pass
    
    def export_via_copypicture(source_obj, output_path: str, width: float, height: float) -> bool:
        """Export an object using CopyPicture method as fallback"""
        try:
            clear_clipboard()
            time.sleep(0.1)
            
            # Copy the object as picture
            source_obj.CopyPicture(Appearance=1, Format=2)  # xlScreen=1, xlBitmap=2
            time.sleep(0.2)
            
            # Create a new chart sheet to paste into (more reliable than ChartObject)
            temp_chart_sheet = workbook.Charts.Add()
            time.sleep(0.1)
            
            try:
                temp_chart_sheet.Paste()
                time.sleep(0.2)
                temp_chart_sheet.Export(output_path, "PNG")
            finally:
                # Clean up temp chart sheet
                excel_app.DisplayAlerts = False
                temp_chart_sheet.Delete()
            
            return os.path.exists(output_path)
            
        except Exception as e:
            print(f"    [CopyPicture Fallback] Failed: {e}")
            return False
    
    try:
        if item_type == 'chartsheet':
            # Chart sheets are straightforward
            print(f"    [ChartSheet] Exporting directly...")
            chart_sheet = workbook.Charts(name)
            chart_sheet.Export(output_path, "PNG")
            
            if not validate_image(output_path):
                print(f"    [ChartSheet] Direct export invalid, trying CopyPicture fallback...")
                if os.path.exists(output_path):
                    os.remove(output_path)
                # Try fallback
                chart_sheet.ChartArea.CopyPicture(Appearance=1, Format=2)
                time.sleep(0.2)
                temp_chart_sheet = workbook.Charts.Add()
                try:
                    temp_chart_sheet.Paste()
                    time.sleep(0.2)
                    temp_chart_sheet.Export(output_path, "PNG")
                finally:
                    excel_app.DisplayAlerts = False
                    temp_chart_sheet.Delete()
                
                if not validate_image(output_path):
                    print(f"    [ChartSheet] Fallback also failed")
                    return False
                    
        else:
            # Worksheet handling
            sheet = workbook.Worksheets(name)
            chart_count = 0
            try:
                chart_count = sheet.ChartObjects().Count
            except:
                pass
            
            print(f"    [Worksheet] Found {chart_count} embedded chart(s)")
            
            if chart_count > 0:
                # Try to export the first chart directly
                chart_obj = sheet.ChartObjects(1)
                print(f"    [Worksheet] Trying direct Chart.Export()...")
                chart_obj.Chart.Export(output_path, "PNG")
                
                if not validate_image(output_path):
                    print(f"    [Worksheet] Direct export failed, trying CopyPicture on ChartObject...")
                    if os.path.exists(output_path):
                        os.remove(output_path)
                    
                    # Fallback: Use CopyPicture on the ChartObject itself
                    success = False
                    for attempt in range(max_retries):
                        try:
                            clear_clipboard()
                            time.sleep(0.1)
                            
                            # Copy the chart object as picture
                            chart_obj.CopyPicture(Appearance=1, Format=2)
                            time.sleep(0.2)
                            
                            # Create temp chart sheet to paste into
                            temp_chart_sheet = workbook.Charts.Add()
                            time.sleep(0.1)
                            
                            try:
                                temp_chart_sheet.Paste()
                                time.sleep(0.2)
                                temp_chart_sheet.Export(output_path, "PNG")
                            finally:
                                excel_app.DisplayAlerts = False
                                temp_chart_sheet.Delete()
                            
                            if validate_image(output_path, min_size=500):
                                success = True
                                print(f"    [Worksheet] CopyPicture fallback succeeded on attempt {attempt + 1}")
                                break
                            else:
                                print(f"    [Worksheet] Attempt {attempt + 1}: CopyPicture fallback validation failed")
                                if os.path.exists(output_path):
                                    os.remove(output_path)
                                    
                        except Exception as e:
                            print(f"    [Worksheet] Attempt {attempt + 1} CopyPicture failed: {e}")
                            time.sleep(0.3)
                    
                    if not success:
                        # Last resort: try copying the entire used range
                        print(f"    [Worksheet] Trying UsedRange CopyPicture as last resort...")
                        try:
                            used_range = sheet.UsedRange
                            if export_via_copypicture(used_range, output_path, used_range.Width, used_range.Height):
                                if validate_image(output_path, min_size=500):
                                    success = True
                                    print(f"    [Worksheet] UsedRange fallback succeeded")
                        except Exception as e:
                            print(f"    [Worksheet] UsedRange fallback failed: {e}")
                    
                    if not success:
                        print(f"    [Worksheet] All methods failed for '{name}'")
                        return False
            else:
                # Worksheet without charts - use CopyPicture with retries
                print(f"    [Worksheet] No charts found, capturing UsedRange...")
                success = False
                last_error = None
                
                for attempt in range(max_retries):
                    try:
                        clear_clipboard()
                        time.sleep(0.1)
                        
                        used_range = sheet.UsedRange
                        
                        # Check if range has content
                        if used_range.Rows.Count == 0 or used_range.Columns.Count == 0:
                            print(f"    [Worksheet] Sheet '{name}' appears empty")
                            return False
                        
                        print(f"    [Worksheet] Attempt {attempt + 1}: UsedRange = {used_range.Rows.Count} rows x {used_range.Columns.Count} cols")
                        
                        # Copy as picture (Appearance=1: xlScreen, Format=2: xlBitmap)
                        used_range.CopyPicture(Appearance=1, Format=2)
                        time.sleep(0.2)
                        
                        # Create temp chart SHEET (more reliable than embedded ChartObject)
                        temp_chart_sheet = workbook.Charts.Add()
                        time.sleep(0.1)
                        
                        try:
                            temp_chart_sheet.Paste()
                            time.sleep(0.2)
                            temp_chart_sheet.Export(output_path, "PNG")
                        finally:
                            excel_app.DisplayAlerts = False
                            temp_chart_sheet.Delete()
                        
                        if validate_image(output_path, min_size=500):
                            success = True
                            break
                        else:
                            print(f"    [Worksheet] Attempt {attempt + 1}: Image validation failed")
                            if os.path.exists(output_path):
                                os.remove(output_path)
                                
                    except Exception as e:
                        last_error = e
                        print(f"    [Worksheet] Attempt {attempt + 1} failed: {e}")
                        time.sleep(0.3)
                        continue
                
                if not success:
                    print(f"    [Worksheet] Failed to capture '{name}' after {max_retries} attempts. Last error: {last_error}")
                    return False
        
        return True
        
    except Exception as e:
        print(f"[Error] Capturing {name}: {e}")
        import traceback
        traceback.print_exc()
        return False


def get_ppt_info(ppt_path: str) -> dict:
    """Get PPT slide info"""
    prs = Presentation(ppt_path)
    slides = []
    for idx, slide in enumerate(prs.slides):
        title = get_slide_title(slide)
        slides.append({
            "page": idx + 1,
            "title": title
        })
    return {
        "total_slides": len(prs.slides),
        "slides": slides,
        "width": prs.slide_width.inches,
        "height": prs.slide_height.inches
    }


def get_slide_title(slide) -> str:
    """Read the title text from a python-pptx slide."""
    try:
        if slide.shapes.title:
            return (slide.shapes.title.text or "").strip()
    except Exception:
        pass
    return ""


def get_ppt_slide_titles(ppt_path: str) -> Dict[int, str]:
    """Return a page -> title mapping for the given PPT."""
    prs = Presentation(ppt_path)
    return {
        idx + 1: get_slide_title(slide)
        for idx, slide in enumerate(prs.slides)
    }


def get_mesh_layout_type(title: str) -> Optional[str]:
    """Return the mesh layout type for supported slide titles."""
    if not title:
        return None

    normalized_title = title.lower()
    if "mesh" not in normalized_title:
        return None
    if "backhaul" in normalized_title:
        return "backhaul"
    if "fronthaul" in normalized_title:
        return "fronthaul"
    return None


def is_mesh_slide_title(title: str) -> bool:
    """True when the slide title matches a supported mesh layout."""
    return get_mesh_layout_type(title) is not None


def get_effective_layout(request: GenerateRequest, slide_title: str) -> Dict[str, float]:
    """Return the placement box for a slide."""
    mesh_layout_type = get_mesh_layout_type(slide_title)
    if mesh_layout_type == "backhaul":
        return dict(MESH_BACKHAUL_LAYOUT)
    if mesh_layout_type == "fronthaul":
        return dict(MESH_FRONTHAUL_LAYOUT)

    return {
        "left": request.img_left,
        "top": request.img_top,
        "width": request.img_width,
        "height": request.img_height,
    }


# ============================================================
# API Endpoints
# ============================================================
@app.get("/", response_class=HTMLResponse)
async def home():
    """Serve the main HTML page"""
    html_path = STATIC_DIR / "index.html"
    if html_path.exists():
        return html_path.read_text(encoding='utf-8')
    return "<html><body><h1>Please add index.html</h1></body></html>"


@app.post("/api/upload-excel")
async def upload_excel(file: UploadFile = File(...)):
    """Upload Excel file and get sheet info"""
    if not file.filename.endswith(('.xlsx', '.xlsm', '.xls')):
        raise HTTPException(400, "必須是 Excel 檔案 (.xlsx, .xlsm, .xls)")
    
    file_id = uuid.uuid4().hex[:8]
    file_path = UPLOAD_DIR / f"{file_id}_{file.filename}"
    
    # Save uploaded file
    with open(file_path, "wb") as f:
        content = await file.read()
        f.write(content)
    
    try:
        info = get_excel_info(str(file_path))
        
        # Store file info
        uploaded_files[file_id] = {
            "type": "excel",
            "path": str(file_path),
            "filename": file.filename
        }
        
        return {
            "status": "success",
            "file_id": file_id,
            "filename": file.filename,
            "worksheets": info["worksheets"],
            "chartsheets": info["chartsheets"]
        }
    except Exception as e:
        # Clean up on error
        if file_path.exists():
            file_path.unlink()
        raise HTTPException(500, f"讀取 Excel 失敗: {str(e)}")


@app.post("/api/upload-ppt")
async def upload_ppt(file: UploadFile = File(...)):
    """Upload PPT template and get slide info"""
    if not file.filename.endswith(('.pptx', '.ppt')):
        raise HTTPException(400, "必須是 PowerPoint 檔案 (.pptx, .ppt)")
    
    file_id = uuid.uuid4().hex[:8]
    file_path = UPLOAD_DIR / f"{file_id}_{file.filename}"
    
    # Save uploaded file
    with open(file_path, "wb") as f:
        content = await file.read()
        f.write(content)
    
    try:
        info = get_ppt_info(str(file_path))
        
        # Store file info
        uploaded_files[file_id] = {
            "type": "ppt",
            "path": str(file_path),
            "filename": file.filename
        }
        
        return {
            "status": "success",
            "file_id": file_id,
            "filename": file.filename,
            **info
        }
    except Exception as e:
        # Clean up on error
        if file_path.exists():
            file_path.unlink()
        raise HTTPException(500, f"讀取 PPT 失敗: {str(e)}")


@app.delete("/api/remove-file/{file_id}")
async def remove_file(file_id: str):
    """Remove an uploaded file"""
    if file_id in uploaded_files:
        file_info = uploaded_files[file_id]
        file_path = Path(file_info["path"])
        if file_path.exists():
            file_path.unlink()
        del uploaded_files[file_id]
        return {"status": "success"}
    raise HTTPException(404, "檔案不存在")


def clear_clipboard():
    """Clear the Windows clipboard"""
    try:
        import win32clipboard
        win32clipboard.OpenClipboard()
        win32clipboard.EmptyClipboard()
        win32clipboard.CloseClipboard()
    except:
        pass


def process_image_mappings(
    mappings: List[ChartMapping],
    prs,
    request: GenerateRequest,
    job_dir: Path,
    slide_titles: Dict[int, str],
) -> List[dict]:
    """Process mappings in image mode - insert charts as static PNG images"""
    import time
    
    results = []
    
    # Group mappings by Excel file
    excel_files = {}
    for m in mappings:
        if m.excel_id not in excel_files:
            excel_files[m.excel_id] = {
                "path": uploaded_files[m.excel_id]["path"],
                "filename": uploaded_files[m.excel_id]["filename"],
                "mappings": []
            }
        excel_files[m.excel_id]["mappings"].append(m)
    
    extracted = {}
    
    # Extract images from each Excel file
    pythoncom.CoInitialize()
    
    excel_app = win32.DispatchEx('Excel.Application')
    excel_app.Visible = False
    excel_app.DisplayAlerts = False
    
    for excel_id, excel_info in excel_files.items():
        print(f"[Image Mode] Opening: {excel_info['filename']}")
        workbook = excel_app.Workbooks.Open(excel_info['path'])
        
        for mapping in excel_info['mappings']:
            key = f"{excel_id}|{mapping.name}"
            safe_name = f"{excel_id}_{mapping.name}"
            for char in '<>:"/\\|?*#':
                safe_name = safe_name.replace(char, '_')
            safe_name = safe_name.replace(' ', '_')
            output_path = str(job_dir / f"{safe_name}.png")
            
            print(f"  Capturing: {mapping.name} (type: {mapping.type})")
            if capture_item(excel_app, workbook, mapping.name, mapping.type, output_path):
                extracted[key] = output_path
                print(f"  [OK] Extracted: {mapping.name} ({os.path.getsize(output_path)} bytes)")
            else:
                print(f"  [FAIL] Failed to extract: {mapping.name}")
        
        workbook.Close(SaveChanges=False)
    
    excel_app.Quit()
    pythoncom.CoUninitialize()
    
    # Insert images into PPT
    for mapping in mappings:
        key = f"{mapping.excel_id}|{mapping.name}"
        excel_filename = uploaded_files[mapping.excel_id]["filename"]
        
        if key not in extracted:
            results.append({
                "name": mapping.name,
                "excel": excel_filename,
                "status": "failed",
                "reason": "擷取失敗"
            })
            continue
        
        slide_idx = mapping.page - 1
        if slide_idx >= len(prs.slides):
            results.append({
                "name": mapping.name,
                "excel": excel_filename,
                "status": "failed",
                "reason": f"第 {mapping.page} 頁不存在"
            })
            continue
        
        slide = prs.slides[slide_idx]
        image_path = extracted[key]
        
        if not os.path.exists(image_path):
            results.append({
                "name": mapping.name,
                "excel": excel_filename,
                "status": "failed",
                "reason": "圖片檔案不存在"
            })
            continue
        
        image_size = os.path.getsize(image_path)
        if image_size < 500:
            results.append({
                "name": mapping.name,
                "excel": excel_filename,
                "status": "failed",
                "reason": f"圖片檔案可能損壞 (大小: {image_size} bytes)"
            })
            continue
        
        try:
            slide_title = slide_titles.get(mapping.page, "")
            layout = get_effective_layout(request, slide_title)
            pic = slide.shapes.add_picture(
                image_path,
                Inches(layout["left"]),
                Inches(layout["top"]),
                width=Inches(layout["width"]),
                height=Inches(layout["height"])
            )
            results.append({
                "name": mapping.name,
                "excel": excel_filename,
                "status": "success",
                "page": mapping.page,
                "mode": "image",
                "mesh_layout": is_mesh_slide_title(slide_title),
            })
        except Exception as img_error:
            print(f"Error adding image to slide: {img_error}")
            results.append({
                "name": mapping.name,
                "excel": excel_filename,
                "status": "failed",
                "reason": f"無法插入圖片: {str(img_error)}"
            })
    
    return results


def process_embedded_mappings(
    mappings: List[ChartMapping],
    ppt_path: str,
    request: GenerateRequest,
    slide_titles: Dict[int, str],
) -> List[dict]:
    """Process mappings in embedded mode - insert charts as editable objects using COM"""
    import time
    
    results = []
    
    # Group mappings by Excel file
    excel_files = {}
    for m in mappings:
        if m.excel_id not in excel_files:
            excel_files[m.excel_id] = {
                "path": uploaded_files[m.excel_id]["path"],
                "filename": uploaded_files[m.excel_id]["filename"],
                "mappings": []
            }
        excel_files[m.excel_id]["mappings"].append(m)
    
    pythoncom.CoInitialize()
    
    excel_app = None
    ppt_app = None
    presentation = None
    
    try:
        # Start Excel - must be visible for clipboard operations
        print("[Embedded Mode] Starting Excel...")
        excel_app = win32.DispatchEx('Excel.Application')
        excel_app.Visible = True
        excel_app.DisplayAlerts = False
        
        # Start PowerPoint - must be visible for paste operations
        print("[Embedded Mode] Starting PowerPoint...")
        ppt_app = win32.DispatchEx('PowerPoint.Application')
        ppt_app.Visible = True
        
        # Open the output PPT
        presentation = ppt_app.Presentations.Open(ppt_path)
        time.sleep(0.5)
        
        # Process each Excel file
        for excel_id, excel_info in excel_files.items():
            print(f"[Embedded Mode] Opening: {excel_info['filename']}")
            workbook = excel_app.Workbooks.Open(os.path.abspath(excel_info['path']))
            time.sleep(0.3)
            
            for mapping in excel_info['mappings']:
                excel_filename = excel_info['filename']
                print(f"  [Embedded Mode] Processing: {mapping.name} -> Page {mapping.page}")
                
                # Check slide exists
                if mapping.page > presentation.Slides.Count:
                    results.append({
                        "name": mapping.name,
                        "excel": excel_filename,
                        "status": "failed",
                        "reason": f"第 {mapping.page} 頁不存在"
                    })
                    continue
                
                try:
                    if mapping.type == 'chartsheet':
                        chart_sheet = workbook.Charts(mapping.name)
                        chart_sheet.Activate()
                        time.sleep(0.2)
                        
                        clear_clipboard()
                        time.sleep(0.1)
                        
                        chart_sheet.ChartArea.Copy()
                    else:
                        sheet = workbook.Worksheets(mapping.name)
                        sheet.Activate()
                        time.sleep(0.2)
                        
                        chart_count = 0
                        try:
                            chart_count = sheet.ChartObjects().Count
                        except:
                            pass
                        
                        if chart_count > 0:
                            chart_obj = sheet.ChartObjects(1)
                            chart_obj.Select()
                            time.sleep(0.2)
                            
                            clear_clipboard()
                            time.sleep(0.1)
                            
                            chart_obj.Chart.ChartArea.Copy()
                        else:
                            print(f"    [Embedded Mode] No chart found, skipping")
                            results.append({
                                "name": mapping.name,
                                "excel": excel_filename,
                                "status": "failed",
                                "reason": "工作表中沒有圖表"
                            })
                            continue
                    
                    time.sleep(0.5)
                    
                    # Paste into PowerPoint
                    slide = presentation.Slides(mapping.page)
                    
                    try:
                        shape = slide.Shapes.Paste()
                        time.sleep(0.3)
                        
                        if hasattr(shape, 'Item'):
                            shape = shape.Item(1)
                        
                        slide_title = slide_titles.get(mapping.page, "")
                        layout = get_effective_layout(request, slide_title)

                        # Convert inches to points (1 inch = 72 points)
                        shape.Left = layout["left"] * 72
                        shape.Top = layout["top"] * 72
                        shape.Width = layout["width"] * 72
                        shape.Height = layout["height"] * 72
                        
                        # Set border to solid black line
                        shape.Line.Visible = -1  # msoTrue = -1 (enable line)
                        shape.Line.ForeColor.RGB = 0  # Black color (RGB 0,0,0)
                        shape.Line.Weight = 0.75  # Line width in points
                        
                        results.append({
                            "name": mapping.name,
                            "excel": excel_filename,
                            "status": "success",
                            "page": mapping.page,
                            "mode": "embedded",
                            "mesh_layout": is_mesh_slide_title(slide_title),
                        })
                        print(f"    [OK] Embedded chart pasted successfully")
                        
                    except Exception as paste_error:
                        print(f"    [FAIL] Paste failed: {paste_error}")
                        results.append({
                            "name": mapping.name,
                            "excel": excel_filename,
                            "status": "failed",
                            "reason": f"貼上失敗: {str(paste_error)}"
                        })
                        
                except Exception as e:
                    print(f"    [ERROR] Error processing {mapping.name}: {e}")
                    results.append({
                        "name": mapping.name,
                        "excel": excel_filename,
                        "status": "failed",
                        "reason": str(e)
                    })
            
            workbook.Close(SaveChanges=False)
        
        # Save and close presentation
        presentation.Save()
        presentation.Close()
        presentation = None
        
    except Exception as e:
        import traceback
        traceback.print_exc()
        raise e
        
    finally:
        if presentation:
            try:
                presentation.Close()
            except:
                pass
        
        if excel_app:
            try:
                excel_app.Quit()
            except:
                pass
        
        if ppt_app:
            try:
                ppt_app.Quit()
            except:
                pass
        
        pythoncom.CoUninitialize()
    
    return results


@app.post("/api/generate")
async def generate_ppt(request: GenerateRequest):
    """Generate PowerPoint with chart mappings
    
    Supports per-mapping chart modes:
    - image: Charts are inserted as static PNG images (default)
    - embedded: Charts are inserted as editable objects
    
    Mixed modes are supported - each mapping can have its own chart_mode.
    """
    
    # Get template path
    if request.template_id not in uploaded_files:
        raise HTTPException(404, "PPT 模板不存在，請重新上傳")
    
    template_path = uploaded_files[request.template_id]["path"]
    
    # Validate mappings
    for m in request.mappings:
        if m.excel_id not in uploaded_files:
            raise HTTPException(404, f"Excel 檔案不存在: {m.excel_id}")
    
    job_id = uuid.uuid4().hex[:8]
    job_dir = OUTPUT_DIR / job_id
    job_dir.mkdir(exist_ok=True)
    
    try:
        slide_titles = get_ppt_slide_titles(template_path)

        # Group mappings by chart_mode
        image_mappings = [m for m in request.mappings if m.chart_mode == "image"]
        embedded_mappings = [m for m in request.mappings if m.chart_mode == "embedded"]
        
        print(f"[Generate] Image mappings: {len(image_mappings)}, Embedded mappings: {len(embedded_mappings)}")
        
        output_filename = f"{request.output_name}.pptx" if not request.output_name.endswith('.pptx') else request.output_name
        output_path = job_dir / output_filename
        
        all_results = []
        
        # Step 1: Process image mappings first (using python-pptx)
        if image_mappings:
            print("[Generate] Processing image mode mappings...")
            prs = Presentation(template_path)
            image_results = process_image_mappings(
                image_mappings,
                prs,
                request,
                job_dir,
                slide_titles,
            )
            all_results.extend(image_results)
            prs.save(str(output_path))
        else:
            # If no image mappings, just copy template
            shutil.copy(template_path, str(output_path))
        
        # Step 2: Process embedded mappings (using COM automation)
        if embedded_mappings:
            print("[Generate] Processing embedded mode mappings...")
            embedded_results = process_embedded_mappings(
                embedded_mappings,
                str(output_path.resolve()),
                request,
                slide_titles,
            )
            all_results.extend(embedded_results)
        
        # Determine mode string for response
        if image_mappings and embedded_mappings:
            mode_str = "mixed"
        elif embedded_mappings:
            mode_str = "embedded"
        else:
            mode_str = "image"
        
        return {
            "status": "success",
            "job_id": job_id,
            "download_url": f"/api/download/{job_id}/{output_filename}",
            "results": all_results,
            "output_file": str(output_path),
            "mode": mode_str
        }
        
    except Exception as e:
        import traceback
        traceback.print_exc()
        raise HTTPException(500, f"產生 PPT 失敗: {str(e)}")


@app.get("/api/download/{job_id}/{filename}")
async def download_file(job_id: str, filename: str):
    """Download generated file"""
    file_path = OUTPUT_DIR / job_id / filename
    
    if not file_path.exists():
        raise HTTPException(404, "檔案不存在或已過期")
    
    return FileResponse(
        file_path,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        filename=filename
    )


if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)
