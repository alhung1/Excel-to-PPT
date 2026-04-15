"""
Netgear NRH Report Generator - Config Version
Just modify the CONFIG section below and run!
"""
import os
import pythoncom
import win32com.client as win32
from pptx import Presentation
from pptx.util import Inches


# ╔══════════════════════════════════════════════════════════════════╗
# ║                    🔧 CONFIGURATION                               ║
# ║               Modify these settings as needed                     ║
# ╚══════════════════════════════════════════════════════════════════╝

# Excel file path
EXCEL_FILE = r"C:\Netgear Projects\NRH testing\RBE770v2\NRH Test Result\RBE773v2 vs RBE773 vs Eero Pro 7_3Pack Test\NRH_RBE773v2 vs RBE773 vs Eero Pro 7_auto channle_Intel BE200_20251016.xlsm"

# PowerPoint template
PPT_TEMPLATE = r"C:\Netgear Projects\NRH_Report.pptx"

# Output file
OUTPUT_PPT = r"C:\Netgear Projects\NRH_Report_Generated.pptx"

# ┌────────────────────────────────────────────────────────────────┐
# │  CHART MAPPINGS                                                 │
# │  Format: ("Chart/Sheet Name", page_number, "type")              │
# │  type: "worksheet" or "chartsheet"                              │
# │  page_number: 1-based (第幾頁)                                   │
# └────────────────────────────────────────────────────────────────┘
CHART_MAPPINGS = [
    ("Metric DUT vs REF#1", 8, "worksheet"),    # 第 8 頁
    ("BI", 9, "chartsheet"),                     # 第 9 頁
    ("DL", 10, "chartsheet"),                    # 第 10 頁
    ("UL", 11, "chartsheet"),                    # 第 11 頁
]

# Image positioning (based on your template)
IMG_LEFT = 0.423      # inches
IMG_TOP = 1.4         # inches
IMG_WIDTH = 12.0      # inches
IMG_HEIGHT = 5.6      # inches

# ╔══════════════════════════════════════════════════════════════════╗
# ║                 END OF CONFIGURATION                              ║
# ╚══════════════════════════════════════════════════════════════════╝

TEMP_DIR = r"C:\Users\alhung\excel-to-ppt\temp_charts"


def capture_item(excel_app, workbook, name, item_type, output_path):
    """Capture a worksheet or chart sheet as image"""
    try:
        if item_type == 'chartsheet':
            chart_sheet = workbook.Charts(name)
            chart_sheet.Export(output_path, "PNG")
        else:
            sheet = workbook.Worksheets(name)
            if sheet.ChartObjects().Count > 0:
                sheet.ChartObjects(1).Chart.Export(output_path, "PNG")
            else:
                used_range = sheet.UsedRange
                used_range.CopyPicture(Appearance=1, Format=2)
                temp_chart = workbook.Charts.Add()
                temp_chart.Paste()
                temp_chart.Export(output_path, "PNG")
                excel_app.DisplayAlerts = False
                temp_chart.Delete()
        return True
    except Exception as e:
        print(f"  ✗ Error: {e}")
        return False


def main():
    print("\n" + "=" * 60)
    print("🔧 Netgear NRH Report Generator")
    print("=" * 60)
    
    print(f"\n📊 Excel: {os.path.basename(EXCEL_FILE)}")
    print(f"📄 Template: {os.path.basename(PPT_TEMPLATE)}")
    print(f"📁 Output: {os.path.basename(OUTPUT_PPT)}")
    
    print("\n📋 Mappings:")
    for name, page, item_type in CHART_MAPPINGS:
        print(f"   {name} → Page {page}")
    
    # Check files
    if not os.path.exists(EXCEL_FILE):
        print(f"\n✗ Excel not found!")
        return
    if not os.path.exists(PPT_TEMPLATE):
        print(f"\n✗ Template not found!")
        return
    
    # Extract from Excel
    print("\n" + "=" * 60)
    print("Step 1: Extracting from Excel")
    print("=" * 60)
    
    os.makedirs(TEMP_DIR, exist_ok=True)
    pythoncom.CoInitialize()
    
    excel_app = win32.DispatchEx('Excel.Application')
    excel_app.Visible = False
    excel_app.DisplayAlerts = False
    workbook = excel_app.Workbooks.Open(EXCEL_FILE)
    
    extracted = {}
    for name, page, item_type in CHART_MAPPINGS:
        safe_name = name.replace(' ', '_').replace('#', '_').replace('/', '_')
        output_path = os.path.join(TEMP_DIR, f"{safe_name}.png")
        
        print(f"\n→ {name}")
        if capture_item(excel_app, workbook, name, item_type, output_path):
            extracted[name] = output_path
            print(f"  ✓ OK")
    
    workbook.Close(SaveChanges=False)
    excel_app.Quit()
    pythoncom.CoUninitialize()
    
    # Insert into PPT
    print("\n" + "=" * 60)
    print("Step 2: Inserting into PowerPoint")
    print("=" * 60)
    
    prs = Presentation(PPT_TEMPLATE)
    
    for name, page, item_type in CHART_MAPPINGS:
        if name not in extracted:
            print(f"\n✗ {name} - skipped")
            continue
        
        slide_idx = page - 1
        if slide_idx >= len(prs.slides):
            print(f"\n✗ Page {page} doesn't exist")
            continue
        
        print(f"\n→ {name} → Page {page}")
        
        slide = prs.slides[slide_idx]
        
        pic = slide.shapes.add_picture(
            extracted[name],
            Inches(IMG_LEFT),
            Inches(IMG_TOP),
            width=Inches(IMG_WIDTH),
            height=Inches(IMG_HEIGHT)
        )
        print(f"  ✓ OK")
    
    prs.save(OUTPUT_PPT)
    
    print("\n" + "=" * 60)
    print(f"🎉 Done! → {OUTPUT_PPT}")
    print("=" * 60)


if __name__ == "__main__":
    main()


